from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.util import Pt
from pptx.enum.chart import XL_LEGEND_POSITION

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = ChartData()
chart_data.categories = ['Q1 Sales', 'Q2 Sales', 'Q3 Sales']
chart_data.add_series('West',    (32.2, 28.4, 34.7))
chart_data.add_series('East',    (24.3, 30.6, 20.2))
chart_data.add_series('Midwest', (20.4, 18.3, 26.2))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

# define axis and labels --------------
chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[0].smooth = True

prs.save('lineChart.pptx')
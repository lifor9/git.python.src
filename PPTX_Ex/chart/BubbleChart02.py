from pptx import Presentation
from pptx.chart.data import XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = BubbleChartData()

series_1 = chart_data.add_series('Series 1')
series_1.add_data_point(0.7, 2.7, 10)
series_1.add_data_point(1.8, 3.2, 4)
series_1.add_data_point(2.6, 0.8, 8)

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
).chart

prs.save('BubbleChart.pptx')
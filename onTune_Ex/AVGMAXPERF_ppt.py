# coding=utf-8
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import  XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
import xlrd
import sys, datetime
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

import psycopg2 as pg2
import os

__version__ = "v1.0.0"
__author__="LIFOR"


def executequery (query) :
    # ================== Database Connection =================== st
    conn_string = "dbname={dbname} user={user} host={host} password={password} port={port}"\
                    .format(dbname='ontune',
                            user='ontune',
                            host='10.50.3.34',
                            password='ontune',
                            port='5432')
    try:
        conn = pg2.connect(database="ontune", user="ontune", password="ontune", host="10.50.3.34", port="5432")
        #conn = psycopg2.connect(conn_string)
    except:
        print("error database connection")

    curs = conn.cursor()
    # ================== Database Connection ===================
    curs.execute(query)

    rows = curs.fetchall()
    # 연결을 종료한다
    curs.close()
    conn.close()

    return rows

hostinfoquer = "SELECT _hostname from hostinfo"
hrows = executequery(hostinfoquer)

print (hrows)
for row in hrows:
   print ("host name : " , row[0])


def text_para_1st(text1):
    # 상단 text box
    left = Cm(1.2); top = Cm(2.0); width = Cm(24); heigh = Cm(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, heigh)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text1
    p.font.size = Pt(12)



def title_para(text1, text2):
    shapes = slide.shapes

    # 라인 1
    shapes = slide.shapes
    left = Inches(0.3)
    top = Inches(0.7)
    width = Inches(9.35)
    height = Inches(0)

    shape = shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, left, top, width, height)

    # 선색
    line = shape.line
    line.color.rgb = RGBColor(0,0,205)
    line.width=Pt(1.5)
    # 채우기
    #fill = shape.fill
    #fill.solid()
    #fill.fore_color.rgb = RGBColor(0,0,160)

    # 이미지 넣기 (그림 1)
    left = Inches(0.3)
    top = Inches(7.15)
    pic2 = slide.shapes.add_picture('image\\ibk1.jpg', left, top)
    # 이미지 넣기 (그림 2)
    left = Inches(9.2)
    top = Inches(7.2)
    pic2 = slide.shapes.add_picture('image\\mtp.jpg', left, top)

    # 제목줄 테스트
    left = Inches(0.3)
    top = Inches(0.3)
    width = Inches(5)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)  # 텍스트 박스 위치 지정
    tf = txBox.text_frame.paragraphs[0]
    tf.text = text1
    tf.font.bold = True
    tf.font.color.rgb = RGBColor(0,0,139)

    # 제목 2
    left = Inches(7.6)
    top = Inches(0.4)
    width = Inches(2)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)  # 텍스트 박스 위치 지정
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text2
    p.font.size = Pt(12)                    # 폰트 크지 지정
    #print(p.font.size.Pt)                   # 폰트 크기 확인
    p.font.color.rgb = RGBColor(53,53,53)   # 글자 색 지정
    p.alignment = PP_ALIGN.RIGHT

    # 라인 2
    shapes = slide.shapes
    left = Inches(0.3)
    top = Inches(7.14)
    width = Inches(9.35)
    height = Inches(0)

    shape = shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, left, top, width, height)

    # 선색
    line = shape.line
    line.color.rgb = RGBColor(0,0,160)
    line.width=Pt(1.0)

    # 채우기
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0,0,160)

# 표지 Page
def print_title_page(base_date):
    title_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = u"시스템 자원 사용 현황"
    subtitle.text = str(base_date)[:4] + u"년" + str(base_date)[4:6]  + u"월"

    #이미지 넣기
    left = Inches(1)
    top = Inches(1)
    pic1 = slide.shapes.add_picture('image\\ibk1.jpg', left, top)
    shapes = slide.shapes

    #선 긋기
    left = Inches(0.7)
    top = Inches(4.82)
    width = Inches(8.6)
    height = Inches(0.06)
    shape = shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE, left, top, width, height)

    #선색
    line = shape.line
    line.color.rgb = RGBColor(0,0,160)
    line.width=Pt(1.0)

# 목차 page
def print_contents_page():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 슬라이드 추가
    title_para(u'목 차 ', '')

    # 텍스트 1
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(4.0)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)  # 텍스트 박스 위치 지정

    tf = txBox.text_frame.paragraphs[0]
    tf.text = """
    01. AP
    02. DB 
    03. 
    04.
    05.
    06.
    07.
    ...
    """
    tf.font.size = Pt(16)                    # 폰트 크지 지정

    # 텍스트 2
    left = Inches(5.0)
    top = Inches(1.0)
    width = Inches(4.0)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)  # 텍스트 박스 위치 지정

    tf = txBox.text_frame.paragraphs[0]
    tf.text = """
    11. AP
    12. DB 
    13. 
    14.
    15.
    16.
    17.
    ...
    """
    tf.font.size = Pt(16)  # 폰트 크지 지정

def print_guide_page():
    global slide 
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 슬라이드 추가
    title_para(u'작성 가이드 ', '')

    # 텍스트
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(4.0)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)  # 텍스트 박스 위치 지정

    tf = txBox.text_frame.paragraphs[0]
    tf.text = """
    1. 시스템별 요약 
    - 업무시간 ... 
    - 최근 ... 
    - 그래프 
    
    2. 일 사용율 그래프 작성 
    - 업무일
    - 각 서버별 
    
    3. 년간 사용률 그래프 
    - 월 단위로 최근 1년간의 
    - 각 서버별 사용률 표시 
    """
    tf.font.size = Pt(12)


# Table Merge Function 정의
def mergeCellsVerically(table, start_row_idx, end_row_idx, col_idx):
    row_count = end_row_idx - start_row_idx + 1
    column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]
    column_cells[0]._tc.set('rowSpan', str(row_count))
    for c in column_cells[1:]:
        c._tc.set('vMerge', '1')

def mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx):
    col_count = end_col_idx - start_col_idx + 1
    row_cells = [c.cells[col_idx] for c in table.rows[row_idx].cells][start_col_idex:end_col_idx]
    row_cells[0]._tc.set('gridSpan', str(col_count))
    for c in row_cells[1:]:
        c._tc.set('hMerge', '1')

def mergeCells(table, start_row_idx, end_row_idx, start_col_idx, end_col_idx):
    for col_idx in range( start_col_idx, end_col_idx + 1):
        mergeCellsVerically(table, start_row_idx, end_row_idx, col_idx)
    for row_idx in range( start_col_idx, end_col_idx + 1):
        mergeCellsHorizontally(table, row_idx, start_col_idx, end_col_idx)

def prt_table2():
    rows = 4
    cols = 10
    left, top, width, height = Cm(1.2), Cm(2.8), Cm(23), Cm(13.5)

def prt_table1():
    rows = 7+3
    cols = 10
    left, top, width, height = Cm(1.2), Cm(2.8), Cm(23), Cm(13.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    col_width =[Cm(3.2), Cm(2.8), Cm(1.9), Cm(2.4), Cm(1.8), Cm(1.8), Cm(1.8), Cm(1.8), Cm(2.0), Cm(3.5)]
    tt = ['A1', 'A2', 'A3', 'A4','A5', 'A6', 'A7', 'A8', 'A9', 'A0']

    for r in range(rows):
        table.rows[r].height = Cm(0.8)

    for c in range(cols):
        table.columns[c].width = col_width[c]
        for r in range(rows):
            table.cell(r, c).text_frame.paragraphs[0].font.size=Pt(9)
            table.cell(r, c).vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        tp = table.cell(0, c).text_frame.paragraphs[0]
        tp.text = tt[c]
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER

    d_tit = table.cell(1, 0).text_frame.paragraphs[0]
    s_tit = table.cell(7, 0).text_frame.paragraphs[0]
    d_tit.text = 'd_tit_text'; d_tit.alignment = PP_ALIGN.CENTER
    s_tit.text = 's_tit_text'; s_tit.alignment = PP_ALIGN.CENTER

    # avgdata = lambda x: round(sum(x) / len(x), 2)

    j = 1
    for i in range(rows -1):
        for qq in range(1, 10):
            table.cell(j, qq).vertival_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        table.cell(j, 1).text_frame.paragraphs[0].text = "table.cell." + str(j).zfill(2)
        table.cell(j, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        j += 1

    return table

def prt_slide_table_p1():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_para('prt_slide_graph_p1', '년  월 시스템 ')
    prt_table1()

def prt_slide_table_p2():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_para('prt_slide_table_p2', '년  월 시스템 ')
    mergeCellsVerically(prt_table1(), 1, 6, 0)

def prt_slide_table_p3():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_para('prt_slide_table_p3', '년  월 시스템 ')
    table = prt_table1()
    # cell = table.cell(1, 0)
    # other_cell = table.cell(6, 0)
    # cell.merge(other_cell)
    table.cell(1, 0).merge(table.cell(6, 0))

def prt_slide_table_p4():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_para('prt_slide_table_p4', '년  월 시스템 ')
    prt_table2()


def prt_slide_graph_p1():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_para('prt_slide_table_p1', '년  월 시스템 ')
    # text_para_1st('월 평균(영업일기준) CPU/MEM 사용율%')
    # 그래프 위 표 그리기
    # prt_table2()

    # 가로 그래프 2개 좌표
    x1, y1, cx1, cy1 = Cm(1.0), Cm(11.0), Cm(11.5), Cm(7.0)
    x2, y2, cx2, cy2 = Cm(13.0), Cm(11.0), Cm(11.5), Cm(7.0)

    # 그래프
    on_label = 0;on_mark = 1
    chart_data1 = ChartData()
    chart_data2 = ChartData()
    # chart_data3 = ChartData()
    # chart_data4 = ChartData()

    # chart_data1.categories = deco_data(data1['data'])
    chart_data1.categories =  ['A', 'B', 'C', 'D']
    chart_data2.categories =  ['A', 'B', 'C', 'D']
    # chart_data1.categories = chart_data2.categories = chart_data3.categories = chart_data4.categories = ['A', 'B', 'C', 'D']

    chart_data1.add_series("CH01", (5,3,7), 2)
    chart_data2.add_series("CH02", (5,3,7), 2)
    # chart_data3.add_series("CH03", (5,3,7), 2)
    # chart_data4.add_series("CH04", (5,3,7), 2)

    add_chart_in_ppt(x1, y1, cx1, cy1, chart_data1)
    add_chart_in_ppt(x2, y2, cx2, cy2, chart_data2)

def prt_slide_graph_p2_1():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_para('prt_slide_graph_p2_1', '년  월 시스템 ')
    text_para_1st('월 평균(영업일기준) CPU/MEM 사용율%')

    # 가로 그래프 1개 좌표
    x1, y1, cx1, cy1 = Cm(1.0), Cm(2.5), Cm(11.5), Cm(7.8)
    # x2, y2, cx2, cy2 = Cm(13.0), Cm(2.5), Cm(11.5), Cm(7.8)
    # x3, y3, cx3, cy3 = Cm(1.0), Cm(10.5), Cm(11.5), Cm(7.8)
    # x4, y4, cx4, cy4 = Cm(13.0), Cm(10.5), Cm(11.5), Cm(7.8)


    # 그래프
    on_label = 0;on_mark = 1
    chart_data1 = ChartData()
    # chart_data2 = ChartData()
    # chart_data3 = ChartData()
    # chart_data4 = ChartData()

    # chart_data1.categories = deco_data(data1['data'])
    chart_data1.categories =  ['A', 'B', 'C', 'D']
    # chart_data1.categories = chart_data2.categories = chart_data3.categories = chart_data4.categories = ['A', 'B', 'C', 'D']

    chart_data1.add_series("CH01", (5,3,7), 2)
    # chart_data2.add_series("CH02", (5,3,7), 2)
    # chart_data3.add_series("CH03", (5,3,7), 2)
    # chart_data4.add_series("CH04", (5,3,7), 2)
    add_chart_in_ppt(x1, y1, cx1, cy1, chart_data1)

def add_chart_in_ppt(x,y,cx,cy, chart_data, on_label=0, on_mark=1):
    if on_mark == 1 :
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x,y,cx,cy, chart_data).chart
    else:
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x,y,cx,cy, chart_data).chart

    # chart.category_axis.tick_labels.font.italic = True
    chart.category_axis.tick_labels.font.size = Pt(8)
    # chart.category_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
    # x 축 라벨 위치
    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    # chart.category_axis.major_unit = 20

    chart.value_axis.visible = True
    # chart.value_axis.has_title = True
    # chart.value_axis.title.text_frame.text = "새로 축 제목"
    # y축 라벨
    chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    chart.value_axis.maximum_scale = 100
    # chart.value_axis.minimum_scale = 100
    chart.value_axis.major_unit = 20
    # 숫자 포맷
    chart.value_axis.tick_labels.number_format = '0"%"'
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    chart.value_axis.has_major_gridlines = True
    # chart.value_axis.has_minor_gridlines = True

    chart.has_legend = True
    chart.has_title = True

    #  표 제목
    p = chart.chart_title.text_frame.paragraphs[0]
    p.text = 'ttttt'
    p.font.size = Pt(10)

    # 범례 라벨 위치
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = True
    chart.legend.font.size = Pt(9)

    if on_label:
        chart.plots[0].has_data_labels = True
        data_labels0 = chart.plots[0].data_labels
        #print(len(chart.series))
        #chart1 = chart.series[1]
        #chart1.plots.data_labels
        data_labels0.number_format = '#.00"*"'
        data_labels0.font.size = Pt(6)
        data_labels0.font.bold = True
        data_labels0.position = XL_TICK_LABEL_POSITION.ABOVE
        # chart.plots[1].has_data_labels = True
        # data_labels1 = chart.plots[1].data_labels
        # data_labels1.number_format = '#.00"%"'
        # data_labels1.font.size = Pt(6)
        # data_labels1.font.bold = True
        # data_labels1.position = XL_TICK_LABEL_POSITION.BELOW

    chart.series.smoot = True



def print_memo_page():
    global slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 슬라이드 추가
    title_para(u'참 고', '')

    # 텍스트
    left = Cm(1.2)
    top = Cm(2.2)
    height = Cm(15.0)
    pic = slide.shapes.add_picture('image\\memo.gif',left, top, height=height)  # 텍스트 박스 위치 지정


def main(base_date):
    global prs
    f1_prefix = "pf_"
    f1_postfix = ".xlsx"

    print(str(base_date))
    outputfile = "TEST3.pptx"
    if os.path.isfile(outputfile):
        print("파일이 존재 : " + os.path.abspath(outputfile))
        os.remove(os.path.abspath(outputfile))

    prs = Presentation()
    print_title_page(base_date)       # 표지 출력
    print_contents_page()   # 목차
    print_guide_page()

    # prt_slide_graph_p1()
    prt_slide_graph_p2_1()

    prt_slide_table_p1()
    prt_slide_table_p2()
    prt_slide_table_p3()

    print_memo_page()       # 참고
    prs.save(outputfile)

if __name__ == '__main__':
    #if len(sys, argv) == 2:
    #    base_date = sys.argv[1]
    #else:
    base_date = input("input query date (yyyymm) : ")

    if base_date == "" :
        bd = datetime.date.today()
        print(bd)
        if len(str(((bd.month -2 ) %12) + 1)) == 1:
            base_date = str(bd.year + ((bd.month -1 ) // 12)) + "0" + str(((bd.month- 2) %12) + 1)
        else:
            base_date = str(bd.year + ((bd.month - 1) // 12)) + str(((bd.month - 2) % 12) + 1)

    print(base_date)

    main(base_date)